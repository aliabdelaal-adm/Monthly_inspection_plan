#!/usr/bin/env python3
"""
مولد قالب التقرير الأسبوعي - Weekly Report Template Generator
يقوم بإنشاء قالب PowerPoint احترافي للتقارير الأسبوعية للتفتيش البيطري

Weekly Report Template Generator
Creates professional PowerPoint templates for weekly veterinary inspection reports
"""

import sys
import io
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ضمان الطباعة بترميز UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

class WeeklyReportGenerator:
    """مولد التقرير الأسبوعي"""
    
    def __init__(self):
        """تهيئة المولد"""
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(13.33)  # 16:9 aspect ratio
        self.presentation.slide_height = Inches(7.5)
        
        # الألوان المهنية
        self.primary_color = RGBColor(35, 54, 160)      # #2336a0
        self.secondary_color = RGBColor(41, 128, 185)   # #2980b9
        self.accent_color = RGBColor(231, 76, 60)       # #e74c3c
        self.text_color = RGBColor(44, 62, 80)          # #2c3e50
        self.light_gray = RGBColor(236, 240, 241)       # #ecf0f1
        
        # أقسام التقرير الرئيسية
        self.report_sections = [
            {
                'title': '1. التفتيش والرقابة البيطرية',
                'subtitle': 'Veterinary Inspection and Monitoring',
                'items': [
                    'عدد التفتيشات الميدانية المنجزة',
                    'نسبة انجاز الخطة (%)',
                    'عدد المحلات/المنشآت التي تمت زيارتها',
                    'مخرجات التفتيش (تعديل فوري، إنذارات، مخالفات)',
                    'التحديات الميدانية'
                ]
            },
            {
                'title': '2. الشكاوى والبلاغات',
                'subtitle': 'Complaints and Reports',
                'items': [
                    'عدد الشكاوى والبلاغات المستلمة',
                    'عدد الشكاوى المغلقة في الوقت المحدد',
                    'عدد الشكاوى العالقة وأسبابها'
                ]
            },
            {
                'title': '3. التوعية والتثقيف',
                'subtitle': 'Awareness and Education',
                'items': [
                    'عدد الفعاليات التوعوية المنفذة',
                    'توعية المحلات (زيارات - برامج)',
                    'حملات مرتبطة بالشكاوى والبلاغات'
                ]
            },
            {
                'title': '4. خدمات الضبط والسيطرة',
                'subtitle': 'Control and Monitoring Services',
                'items': [
                    'عدد الحيوانات التي تم التعامل معها',
                    'نسبة إغلاق الحالات في الوقت المحدد',
                    'أنواع الحيوانات (قطط، كلاب، حيوانات أخرى)',
                    'مخرجات الضبط (إيواء، تعقيم، تبني...)'
                ]
            },
            {
                'title': '5. الزيارات والمتابعة',
                'subtitle': 'Visits and Follow-up',
                'items': [
                    'عدد الزيارات على الملاجئ ودور الإيواء',
                    'الملاحظات أو التحديات خلال الزيارات'
                ]
            },
            {
                'title': '6. العقود والمتابعة',
                'subtitle': 'Contracts and Follow-up',
                'items': [
                    'متابعة عقود التشغيل (الاجتماعات الأسبوعية ومخرجاتها)',
                    'أي ملاحظات أو تحديثات من المقاولين/المشغلين'
                ]
            },
            {
                'title': '7. التطوير والجودة',
                'subtitle': 'Development and Quality',
                'items': [
                    'مقترحات تحسين اللوائح والسياسات',
                    'متابعات الجودة والمخاطر',
                    'أي برامج تدريبية أو خطط تصحيحية مقترحة'
                ]
            }
        ]
    
    def calculate_week_dates(self, week_number, year=2025):
        """حساب تواريخ الأسبوع (الاثنين إلى الأحد)"""
        # أول يوم اثنين في السنة
        jan_1 = datetime(year, 1, 1)
        days_to_monday = (7 - jan_1.weekday()) % 7
        if jan_1.weekday() == 0:  # إذا كان أول يناير يوم الاثنين
            days_to_monday = 0
        else:
            days_to_monday = 7 - jan_1.weekday()
        
        first_monday = jan_1 + timedelta(days=days_to_monday)
        
        # حساب الاثنين للأسبوع المطلوب
        week_monday = first_monday + timedelta(weeks=week_number-1)
        week_sunday = week_monday + timedelta(days=6)
        submission_tuesday = week_sunday + timedelta(days=2)
        
        return week_monday, week_sunday, submission_tuesday
    
    def add_title_slide(self, week_number, week_start, week_end, submission_date):
        """إضافة شريحة العنوان الرئيسية"""
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # خلفية متدرجة
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.light_gray
        
        # العنوان الرئيسي
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = f"التقرير الأسبوعي رقم {week_number}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.color.rgb = self.primary_color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # العنوان الفرعي بالإنجليزية
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.33), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Weekly Report No. {week_number}"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = self.secondary_color
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # فترة التقرير
        period_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9.33), Inches(1))
        period_frame = period_box.text_frame
        period_frame.text = f"من يوم الاثنين {week_start.strftime('%d/%m/%Y')} الى يوم الأحد {week_end.strftime('%d/%m/%Y')}"
        period_para = period_frame.paragraphs[0]
        period_para.font.size = Pt(24)
        period_para.font.color.rgb = self.text_color
        period_para.alignment = PP_ALIGN.CENTER
        
        # تاريخ التسليم
        submission_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9.33), Inches(0.8))
        submission_frame = submission_box.text_frame
        submission_frame.text = f"يرفع يوم الثلاثاء {submission_date.strftime('%d/%m/%Y')}"
        submission_para = submission_frame.paragraphs[0]
        submission_para.font.size = Pt(20)
        submission_para.font.color.rgb = self.accent_color
        submission_para.font.bold = True
        submission_para.alignment = PP_ALIGN.CENTER
        
        # شعار أو رمز (اختياري)
        logo_box = slide.shapes.add_textbox(Inches(5.5), Inches(6), Inches(2.33), Inches(0.8))
        logo_frame = logo_box.text_frame
        logo_frame.text = "🏥 التفتيش البيطري"
        logo_para = logo_frame.paragraphs[0]
        logo_para.font.size = Pt(18)
        logo_para.font.color.rgb = self.primary_color
        logo_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_section_slide(self, section):
        """إضافة شريحة لقسم من أقسام التقرير"""
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # خلفية
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # شريط عنوان ملون
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            Inches(13.33), Inches(1.2)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.primary_color
        header_shape.line.fill.background()
        
        # عنوان القسم
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = section['title']
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.RIGHT
        
        # العنوان الفرعي بالإنجليزية
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.33), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = section['subtitle']
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.RIGHT
        
        # منطقة المحتوى الرئيسي
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.text = "النقاط الرئيسية:"
        
        # إضافة النقاط
        for item in section['items']:
            p = content_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_color
            p.space_after = Pt(6)
        
        # منطقة للصور والرسوم البيانية
        chart_area = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(1.5), 
            Inches(4.5), Inches(2.5)
        )
        chart_fill = chart_area.fill
        chart_fill.solid()
        chart_fill.fore_color.rgb = self.light_gray
        chart_area.line.color.rgb = self.secondary_color
        chart_area.line.width = Pt(2)
        
        # نص للصور والرسوم
        chart_text_box = slide.shapes.add_textbox(Inches(8.7), Inches(2.2), Inches(4.1), Inches(1.1))
        chart_text_frame = chart_text_box.text_frame
        chart_text_frame.text = "منطقة للصور والرسوم البيانية\nArea for Images & Charts"
        chart_text_para = chart_text_frame.paragraphs[0]
        chart_text_para.font.size = Pt(14)
        chart_text_para.font.color.rgb = self.secondary_color
        chart_text_para.alignment = PP_ALIGN.CENTER
        
        # منطقة للتحاليل الإحصائية
        stats_area = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(4.2), 
            Inches(4.5), Inches(2.5)
        )
        stats_fill = stats_area.fill
        stats_fill.solid()
        stats_fill.fore_color.rgb = self.light_gray
        stats_area.line.color.rgb = self.accent_color
        stats_area.line.width = Pt(2)
        
        # نص للتحاليل الإحصائية
        stats_text_box = slide.shapes.add_textbox(Inches(8.7), Inches(4.9), Inches(4.1), Inches(1.1))
        stats_text_frame = stats_text_box.text_frame
        stats_text_frame.text = "منطقة للتحاليل الإحصائية\nStatistical Analysis Area"
        stats_text_para = stats_text_frame.paragraphs[0]
        stats_text_para.font.size = Pt(14)
        stats_text_para.font.color.rgb = self.accent_color
        stats_text_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_summary_slide(self):
        """إضافة شريحة الملخص والخلاصة"""
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # خلفية
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # شريط العنوان
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            Inches(13.33), Inches(1.2)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = self.accent_color
        header_shape.line.fill.background()
        
        # عنوان الملخص
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "الملخص والخلاصة | Summary & Conclusion"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # منطقة الملخص الرئيسي
        summary_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(2))
        summary_frame = summary_box.text_frame
        summary_frame.text = "النقاط الرئيسية للأسبوع:\n\n• الإنجازات الرئيسية\n• التحديات المواجهة\n• الخطط المستقبلية"
        summary_para = summary_frame.paragraphs[0]
        summary_para.font.size = Pt(20)
        summary_para.font.color.rgb = self.text_color
        
        # منطقة للإحصائيات الإجمالية
        stats_total_area = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), 
            Inches(11.33), Inches(2.5)
        )
        stats_total_fill = stats_total_area.fill
        stats_total_fill.solid()
        stats_total_fill.fore_color.rgb = self.light_gray
        stats_total_area.line.color.rgb = self.primary_color
        stats_total_area.line.width = Pt(3)
        
        # نص الإحصائيات الإجمالية
        stats_total_text_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.33), Inches(1.5))
        stats_total_text_frame = stats_total_text_box.text_frame
        stats_total_text_frame.text = "الإحصائيات الإجمالية للأسبوع\nWeekly Total Statistics\n\n[منطقة لعرض الإحصائيات والأرقام الرئيسية]"
        stats_total_text_para = stats_total_text_frame.paragraphs[0]
        stats_total_text_para.font.size = Pt(18)
        stats_total_text_para.font.color.rgb = self.primary_color
        stats_total_text_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def generate_template(self, week_number=1, year=2025):
        """إنشاء القالب الكامل"""
        print(f"🚀 بدء إنشاء قالب التقرير الأسبوعي رقم {week_number} لسنة {year}")
        
        # حساب التواريخ
        week_start, week_end, submission_date = self.calculate_week_dates(week_number, year)
        
        print(f"📅 فترة التقرير: من {week_start.strftime('%d/%m/%Y')} إلى {week_end.strftime('%d/%m/%Y')}")
        print(f"📤 تاريخ التسليم: {submission_date.strftime('%d/%m/%Y')}")
        
        # إضافة الشرائح
        print("📄 إضافة شريحة العنوان...")
        self.add_title_slide(week_number, week_start, week_end, submission_date)
        
        print("📋 إضافة شرائح الأقسام...")
        for i, section in enumerate(self.report_sections, 1):
            print(f"   {i}/7 - {section['title']}")
            self.add_section_slide(section)
        
        print("📊 إضافة شريحة الملخص...")
        self.add_summary_slide()
        
        # حفظ الملف
        filename = f"التقرير_الأسبوعي_رقم_{week_number}_{year}.pptx"
        self.presentation.save(filename)
        
        print(f"✅ تم إنشاء القالب بنجاح: {filename}")
        print(f"📈 إجمالي الشرائح: {len(self.presentation.slides)}")
        
        return filename

def main():
    """الدالة الرئيسية"""
    print("=== مولد قالب التقرير الأسبوعي ===")
    print("Weekly Report Template Generator")
    print()
    
    try:
        # إنشاء مولد التقرير
        generator = WeeklyReportGenerator()
        
        # إنشاء عدة قوالب كأمثلة
        examples = [
            (1, 2025),  # التقرير الأسبوعي رقم 1
            (2, 2025),  # التقرير الأسبوعي رقم 2
            (3, 2025),  # التقرير الأسبوعي رقم 3
        ]
        
        print("🎯 إنشاء قوالب أمثلة:")
        for week_num, year in examples:
            filename = generator.generate_template(week_num, year)
            print(f"   ✓ {filename}")
            print()
        
        print("🎉 تم إنشاء جميع القوالب بنجاح!")
        print()
        print("📝 ملاحظات الاستخدام:")
        print("   • كل قالب يحتوي على 9 شرائح (عنوان + 7 أقسام + ملخص)")
        print("   • توجد أماكن مخصصة للصور والرسوم البيانية والتحاليل الإحصائية")
        print("   • التواريخ محسوبة تلقائياً (الاثنين إلى الأحد)")
        print("   • تاريخ التسليم محدد يوم الثلاثاء التالي")
        
    except Exception as e:
        print(f"❌ خطأ في إنشاء القالب: {e}")
        return False
    
    return True

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
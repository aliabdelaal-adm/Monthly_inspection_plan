#!/usr/bin/env python3
"""
مولد التقرير الأسبوعي مع البيانات - Weekly Report Generator with Data Integration
يقوم بإنشاء تقارير أسبوعية مع دمج البيانات من ملف plan-data.json

Weekly Report Generator with Data Integration
Creates weekly reports with data integration from plan-data.json
"""

import sys
import json
import argparse
from datetime import datetime, timedelta
from collections import Counter, defaultdict
from weekly_report_generator import WeeklyReportGenerator
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class WeeklyReportWithData(WeeklyReportGenerator):
    """مولد التقرير الأسبوعي مع البيانات"""
    
    def __init__(self, data_file='plan-data.json'):
        """تهيئة المولد مع ملف البيانات"""
        super().__init__()
        self.data_file = data_file
        self.inspection_data = self.load_inspection_data()
    
    def load_inspection_data(self):
        """تحميل بيانات التفتيش من الملف"""
        try:
            with open(self.data_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
            print(f"✅ تم تحميل بيانات التفتيش: {len(data.get('inspectionData', []))} مدخل")
            return data
        except FileNotFoundError:
            print(f"⚠️ لم يتم العثور على ملف البيانات: {self.data_file}")
            return {"inspectionData": [], "inspectors": [], "areas": [], "shops": []}
        except Exception as e:
            print(f"❌ خطأ في تحميل البيانات: {e}")
            return {"inspectionData": [], "inspectors": [], "areas": [], "shops": []}
    
    def filter_data_by_week(self, start_date, end_date):
        """تصفية البيانات حسب الأسبوع المحدد"""
        week_data = []
        
        for entry in self.inspection_data.get('inspectionData', []):
            try:
                entry_date = datetime.strptime(entry['day'], '%Y-%m-%d')
                if start_date <= entry_date <= end_date:
                    week_data.append(entry)
            except (KeyError, ValueError):
                continue
        
        return week_data
    
    def calculate_weekly_statistics(self, week_data):
        """حساب الإحصائيات الأسبوعية"""
        if not week_data:
            return {
                'total_inspections': 0,
                'total_shops': 0,
                'total_inspectors': 0,
                'areas_covered': 0,
                'inspector_breakdown': {},
                'area_breakdown': {},
                'shift_breakdown': {},
                'daily_breakdown': {}
            }
        
        # حساب الإحصائيات الأساسية
        total_inspections = len(week_data)
        
        all_shops = []
        inspectors = set()
        areas = set()
        shifts = []
        days = []
        
        inspector_counts = Counter()
        area_counts = Counter()
        shift_counts = Counter()
        daily_counts = Counter()
        
        for entry in week_data:
            inspectors.add(entry.get('inspector', 'غير محدد'))
            areas.add(entry.get('area', 'غير محدد'))
            shifts.append(entry.get('shift', 'غير محدد'))
            
            # تحويل التاريخ لعرض أفضل
            try:
                date_obj = datetime.strptime(entry['day'], '%Y-%m-%d')
                day_name = date_obj.strftime('%A')  # اسم اليوم بالإنجليزية
                days.append(f"{entry['day']} ({day_name})")
            except:
                days.append(entry.get('day', 'غير محدد'))
            
            # عد المحلات
            shops_in_entry = entry.get('shops', [])
            all_shops.extend(shops_in_entry)
            
            # إحصائيات تفصيلية
            inspector_counts[entry.get('inspector', 'غير محدد')] += 1
            area_counts[entry.get('area', 'غير محدد')] += 1
            shift_counts[entry.get('shift', 'غير محدد')] += 1
            daily_counts[entry.get('day', 'غير محدد')] += 1
        
        return {
            'total_inspections': total_inspections,
            'total_shops': len(all_shops),
            'unique_shops': len(set(all_shops)),
            'total_inspectors': len(inspectors),
            'areas_covered': len(areas),
            'inspector_breakdown': dict(inspector_counts),
            'area_breakdown': dict(area_counts),
            'shift_breakdown': dict(shift_counts),
            'daily_breakdown': dict(daily_counts),
            'inspectors_list': list(inspectors),
            'areas_list': list(areas),
        }
    
    def add_data_enriched_summary_slide(self, statistics):
        """إضافة شريحة الملخص مع البيانات المُثراة"""
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # خلفية
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.light_gray
        
        # عنوان الملخص
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "📊 الملخص التنفيذي والإحصائيات الأسبوعية"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.color.rgb = self.primary_color
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # الإحصائيات الرئيسية
        stats_text = f"""
📈 الإحصائيات الرئيسية:
• إجمالي عمليات التفتيش: {statistics['total_inspections']}
• إجمالي المحلات المُفتشة: {statistics['total_shops']} ({statistics['unique_shops']} محل فريد)
• عدد المفتشين المشاركين: {statistics['total_inspectors']}
• عدد المناطق المغطاة: {statistics['areas_covered']}

👥 توزيع المفتشين:"""
        
        # إضافة توزيع المفتشين
        for inspector, count in statistics['inspector_breakdown'].items():
            stats_text += f"\n• {inspector}: {count} تفتيش"
        
        stats_text += "\n\n🏘️ توزيع المناطق:"
        
        # إضافة توزيع المناطق
        for area, count in statistics['area_breakdown'].items():
            stats_text += f"\n• {area}: {count} تفتيش"
        
        # منطقة الإحصائيات
        stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.5))
        stats_frame = stats_box.text_frame
        stats_frame.text = stats_text
        stats_para = stats_frame.paragraphs[0]
        stats_para.font.size = Pt(14)
        stats_para.font.color.rgb = self.text_color
        
        return slide
    
    def generate_data_driven_template(self, start_date, end_date, week_number=None):
        """إنشاء قالب مدفوع بالبيانات"""
        # تصفية البيانات للأسبوع المحدد
        week_data = self.filter_data_by_week(start_date, end_date)
        
        # حساب الإحصائيات
        statistics = self.calculate_weekly_statistics(week_data)
        
        print(f"📊 تم العثور على {len(week_data)} عملية تفتيش للفترة المحددة")
        print(f"📈 إحصائيات مُحسبة: {statistics['total_inspections']} تفتيش، {statistics['unique_shops']} محل")
        
        if week_number is None:
            # حساب رقم الأسبوع
            year = start_date.year
            jan_1 = datetime(year, 1, 1)
            days_to_monday = (7 - jan_1.weekday()) % 7 if jan_1.weekday() != 0 else 0
            first_monday = jan_1 + timedelta(days=days_to_monday)
            week_number = ((start_date - first_monday).days // 7) + 1
        
        submission_date = end_date + timedelta(days=2)
        
        print(f"🚀 بدء إنشاء التقرير الأسبوعي المُثرى بالبيانات")
        print(f"📅 فترة التقرير: من {start_date.strftime('%d/%m/%Y')} إلى {end_date.strftime('%d/%m/%Y')}")
        print(f"📤 تاريخ التسليم: {submission_date.strftime('%d/%m/%Y')}")
        
        # إضافة الشرائح
        print("📄 إضافة شريحة العنوان...")
        self.add_title_slide(week_number, start_date, end_date, submission_date)
        
        print("📋 إضافة شرائح الأقسام...")
        for i, section in enumerate(self.report_sections, 1):
            print(f"   {i}/7 - {section['title']}")
            self.add_section_slide(section)
        
        print("📊 إضافة شريحة الملخص المُثراة بالبيانات...")
        self.add_data_enriched_summary_slide(statistics)
        
        # حفظ الملف
        filename = f"التقرير_الأسبوعي_رقم_{week_number}_{start_date.year}_مع_البيانات.pptx"
        self.presentation.save(filename)
        
        print(f"✅ تم إنشاء القالب المُثرى بنجاح: {filename}")
        print(f"📈 إجمالي الشرائح: {len(self.presentation.slides)}")
        
        # طباعة ملخص الإحصائيات
        print("\n📊 ملخص البيانات المُدمجة:")
        print(f"   📝 عمليات التفتيش: {statistics['total_inspections']}")
        print(f"   🏪 المحلات: {statistics['total_shops']} (فريد: {statistics['unique_shops']})")
        print(f"   👥 المفتشين: {statistics['total_inspectors']}")
        print(f"   🏘️ المناطق: {statistics['areas_covered']}")
        
        return filename

def main():
    """الدالة الرئيسية"""
    parser = argparse.ArgumentParser(
        description="مولد التقرير الأسبوعي مع البيانات - Weekly Report Generator with Data",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
أمثلة على الاستخدام:
Examples:

# إنشاء تقرير مُثرى بالبيانات
python3 weekly_report_with_data.py --start 22/09/2025 --end 28/09/2025

# تحديد ملف بيانات مختلف
python3 weekly_report_with_data.py --start 22/09/2025 --end 28/09/2025 --data custom-data.json
        """
    )
    
    parser.add_argument(
        '--start', '-s',
        required=True,
        help='تاريخ بداية التقرير (الاثنين) - تنسيق: DD/MM/YYYY'
    )
    
    parser.add_argument(
        '--end', '-e',
        required=True,
        help='تاريخ نهاية التقرير (الأحد) - تنسيق: DD/MM/YYYY'
    )
    
    parser.add_argument(
        '--week', '-w',
        type=int,
        help='رقم الأسبوع (اختياري - سيتم حسابه تلقائياً إذا لم يُحدد)'
    )
    
    parser.add_argument(
        '--data', '-d',
        default='plan-data.json',
        help='ملف البيانات (افتراضي: plan-data.json)'
    )
    
    args = parser.parse_args()
    
    try:
        # تحويل التواريخ
        start_date = datetime.strptime(args.start, "%d/%m/%Y")
        end_date = datetime.strptime(args.end, "%d/%m/%Y")
        
        # التحقق من صحة التواريخ
        if (end_date - start_date).days != 6:
            print("❌ يجب أن تكون فترة التقرير 7 أيام بالضبط")
            return False
        
        if start_date.weekday() != 0:
            print("❌ يجب أن يبدأ التقرير بيوم الاثنين")
            return False
        
        # إنشاء التقرير
        print("=== مولد التقرير الأسبوعي مع البيانات ===")
        print("Weekly Report Generator with Data Integration")
        print()
        
        generator = WeeklyReportWithData(args.data)
        filename = generator.generate_data_driven_template(start_date, end_date, args.week)
        
        print()
        print("🎉 تم إنشاء التقرير المُثرى بالبيانات بنجاح!")
        print(f"📁 اسم الملف: {filename}")
        
        return True
        
    except ValueError as e:
        print(f"❌ خطأ في التاريخ: {e}")
        return False
    except Exception as e:
        print(f"❌ خطأ في إنشاء التقرير: {e}")
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)